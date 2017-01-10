#!/usr/bin/env python
# -*- coding: utf-8 -*-

from pptx import *
import os
from datetime import *
import gc
from CSV_Manager import CSV_Manager

class Ppt_creator(object):

    def __init__(self, configuration_file, layout_identifier, data_path, parameter):
        self.layout_identifier = layout_identifier
        self.models = configuration_file.MODELS
        self.presentation_template_file = configuration_file.PPT_TEMPLATE
        self.prediction_start_date = configuration_file.TODAY
        self.data_path = data_path
        self.presentation = Presentation(self.presentation_template_file)
        self.slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[self.layout_identifier])
        self.params = configuration_file.FULL_PARAMS
        self.legends_dictionary = configuration_file.PARAMETER_LEGENDS
        self.param = parameter
        self.week_days = configuration_file.WEEK_DAYS
        self.csv_file = configuration_file.RELIABILITY_DATA_FILE

        self.files = self.load_files(self.data_path)
        self.gfs = self.get_model_files('gfs', self.files)
        self.ifs = self.get_model_files('ifs', self.files)
        self.cmc = self.get_model_files('cmc', self.files)
        self.gum = self.get_model_files('gum', self.files)
        self.arpege = self.get_model_files('arpege', self.files)

        self.picture_placeholders_identifiers = self.get_placeholder_identifiers('Picture', self.slide)
        self.legend_placeholder_identifier = []
        self.legend_placeholder_identifier.append(self.picture_placeholders_identifiers[0])
        self.picture_placeholders_identifiers = self.picture_placeholders_identifiers[1:]
        self.text_placeholder_identifiers = self.get_placeholder_identifiers('Text', self.slide)
        self.dates_placeholders_identifiers =  self.text_placeholder_identifiers[0:3]
        self.prediction_release_date = []
        self.prediction_release_date.append(self.text_placeholder_identifiers[3])
        self.model_reliability_percent_identifiers = self.text_placeholder_identifiers[4:-3]
        self.day_in_text_form_identifiers = self.text_placeholder_identifiers[-3:]
        csv_data = CSV_Manager(self.csv_file).get_csv_data_list()
        self.reliability_percent = csv_data

        self.create_presentation()

    def create_presentation(self):
        self.fill_layouts(self.layout_identifier)

        # Clear Garbage Collector to release memory of internal python paths arrays
        gc.collect()

    def load_files(self, data_path):
        files = []
        for file in os.listdir(data_path):
            if '.png' in file:
                files.append(data_path + file)
        return files

    def get_model_files(self, model_name, files_list):
        files = []
        for file in files_list:
            if model_name in file:
                files.append(file)
        return files

    def get_shift_number_from_layout_identifier(self, layout_identifier):
        if layout_identifier == 0:
            return 0
        elif layout_identifier == 1:
            return 3
        elif layout_identifier == 2:
            return 6
        else:
            print "Layout identifier: " + str(layout_identifier) + " is not valid - insert just 0,1 or 2"

    def fill_dates_to_layout(self, shift_number, dates_placeholders_identifiers):
        i = 0
        for identifier in dates_placeholders_identifiers:
            date = self.prediction_start_date + timedelta(days=shift_number + 1) + timedelta(days=i)
            text_shape = self.slide.placeholders[identifier]
            text_shape.text = str(date.day) + '. ' + str(date.month) + '. '
            i += 1

    def fill_layout_images(self, shift_number):
        i = 0

        w, q, x, y, z = 0, 0, 0, 0, 0
        for identifier in self.picture_placeholders_identifiers:
            placeholder = self.slide.placeholders[identifier]
            if i == 0 or i == 5 or i == 10:
                placeholder.insert_picture(self.ifs[shift_number + w])
                w += 1
            elif i == 1 or i == 6 or i == 11:
                placeholder.insert_picture(self.gfs[shift_number + q])
                q += 1
            elif i == 2 or i == 7 or i == 12:
                placeholder.insert_picture(self.cmc[shift_number + x])
                x += 1
            elif i == 3 or i == 8 or i == 13:
                placeholder.insert_picture(self.gum[shift_number + y])
                y += 1
            elif i == 4 or i == 9 or i == 14:
                placeholder.insert_picture(self.arpege[shift_number + z])
                z += 1
            i += 1

        self.fill_prediction_release_date(self.prediction_release_date)
        self.fill_dates_to_layout(shift_number, self.dates_placeholders_identifiers)
        self.fill_text_day_values(self.day_in_text_form_identifiers, shift_number)
        legend_image_file = self.get_legend()
        self.fill_legend(legend_image_file, self.legend_placeholder_identifier)
        self.fill_reliability_values(self.model_reliability_percent_identifiers)
        self.save_pres_to_template(self.presentation, self.presentation_template_file)

    def fill_layouts(self, layout_identifier):
        shift_number = self.get_shift_number_from_layout_identifier(layout_identifier)
        self.fill_layout_images(shift_number)

    def save_pres_to_template(self, pres_object, template):
        try:
            pres_object.save(template)
        except Exception as e:
            print "Saving presentation failed with exception: " + str(e)

    def get_placeholder_identifiers(self, placeholder_type, slide):
        placeholders_identifiers = []
        for shape in slide.shapes:
            if placeholder_type in shape.name:
                placeholders_identifiers.append(shape.placeholder_format.idx)
        return placeholders_identifiers

    def get_legend(self):
        if self.param.replace('_', '').upper() in self.data_path:
            legend_path = self.legends_dictionary.get(self.param)
            return legend_path
        else:
            print "Legend image unreachable!"

    def fill_legend(self, legend_file, legend_placeholder):
        for identifier in legend_placeholder:
            placeholder = self.slide.placeholders[identifier]
            placeholder.insert_picture(legend_file)

    def fill_prediction_release_date(self, prediction_release_date_identifiers):
        i = 0
        for identifier in prediction_release_date_identifiers:
            date = self.prediction_start_date
            text_shape = self.slide.placeholders[identifier]
            text_shape.text = str(date.day) + '. ' + str(date.month) + '. ' + str(date.year)
            i += 1

    def fill_reliability_values(self, reliability_text_identifiers):
        i = 0
        for identifier in reliability_text_identifiers:
            reliability_text = self.reliability_percent[self.layout_identifier][i]
            text_shape = self.slide.placeholders[identifier]
            text_shape.text = reliability_text
            i += 1

    def fill_text_day_values(self, day_in_text_form_identifiers, shift_number):
        i = 0
        for identifier in day_in_text_form_identifiers:
            date = self.prediction_start_date + timedelta(days=shift_number + 1) + timedelta(days=i)
            text_shape = self.slide.placeholders[identifier]
            week_day = self.week_days.get(date.weekday())
            text_shape.text = str(week_day)
            i += 1

